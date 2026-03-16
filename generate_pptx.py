"""
Generate a PowerPoint presentation from the Workshop GitHub + Azure DevOps website content.
Run: python3 generate_pptx.py
Output: workshop_github_azure_devops.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours (match index.css) ──────────────────────────────────────────
DARK       = RGBColor(0x1F, 0x2A, 0x35)   # --text
MUTED      = RGBColor(0x56, 0x61, 0x6F)   # --muted
ACCENT     = RGBColor(0xF3, 0x6A, 0x31)   # --accent  (orange)
SECONDARY  = RGBColor(0x0F, 0x7C, 0x8F)   # --secondary (teal)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xFD, 0xF7, 0xEE)   # --bg
TEAL_DARK  = RGBColor(0x0F, 0x51, 0x60)   # cta gradient end


def set_bg(slide, color: RGBColor):
    """Set slide background to a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=None, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"\u25cf  {item}"
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ── Slide helpers ─────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def slide_title(prs):
    """Slide 1 – Title / Hero"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_LIGHT)

    # Teal accent bar on the left
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, SECONDARY)

    # Orange top strip
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    # Eyebrow
    add_textbox(slide, "WORKSHOP GRATUITO + CERTIFICADO",
                Inches(0.5), Inches(0.6), Inches(8), Inches(0.5),
                font_size=11, bold=True, color=SECONDARY)

    # Title
    add_textbox(slide,
                "GitHub + Azure DevOps para automatizar builds, testes e deploy",
                Inches(0.5), Inches(1.2), Inches(8), Inches(2),
                font_size=36, bold=True, color=DARK)

    # Sub-title paragraph
    add_textbox(slide,
                "Primeiros passos para estruturar pipelines de CI/CD com clareza,\n"
                "reduzindo tarefas manuais e acelerando entregas.",
                Inches(0.5), Inches(3.3), Inches(7.5), Inches(1),
                font_size=16, color=MUTED)

    # Highlighted quote box (teal)
    add_rect(slide, Inches(0.5), Inches(4.4), Inches(7.5), Inches(0.85),
             RGBColor(0xD5, 0xF0, 0xF5))
    add_textbox(slide,
                "Em apenas 1 hora, veja como estruturar um pipeline real de build, "
                "testes e deploy usando GitHub e Azure DevOps.",
                Inches(0.7), Inches(4.45), Inches(7.2), Inches(0.8),
                font_size=13, color=RGBColor(0x21, 0x4A, 0x55), italic=True)

    # CTA button area (orange rounded rect visual)
    btn = add_rect(slide, Inches(0.5), Inches(5.5), Inches(3.5), Inches(0.55),
                   ACCENT)
    add_textbox(slide, "Garantir minha vaga gratuitamente",
                Inches(0.55), Inches(5.53), Inches(3.4), Inches(0.5),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right-side teal panel
    add_rect(slide, Inches(9), Inches(0.3), Inches(3.9), Inches(6.5),
             SECONDARY)
    add_textbox(slide, "Workshop\nGitHub +\nAzure DevOps",
                Inches(9.2), Inches(1.5), Inches(3.5), Inches(3),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "CI/CD · Automação · DevOps",
                Inches(9.2), Inches(4.6), Inches(3.5), Inches(0.6),
                font_size=12, color=RGBColor(0xD5, 0xF0, 0xF5),
                align=PP_ALIGN.CENTER)


def slide_audience(prs):
    """Slide 2 – Para quem é este workshop?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "Para quem é este workshop?",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=DARK)
    add_textbox(slide,
                "Este workshop é ideal para profissionais e equipes que querem elevar "
                "o nível de automação no ciclo de desenvolvimento.",
                Inches(0.5), Inches(1.1), Inches(12), Inches(0.6),
                font_size=14, color=MUTED)

    cards = [
        ("Desenvolvedores",
         "Querem aprender a automatizar build e deploy com ferramentas usadas no mercado."),
        ("Engenheiros de software",
         "Possuem interesse em DevOps e CI/CD para tornar o processo de entrega mais confiável."),
        ("Iniciantes em automação",
         "Querem conhecer GitHub Actions e Azure DevOps de forma objetiva e aplicada."),
        ("Times de tecnologia",
         "Buscam melhorar o fluxo de desenvolvimento, testes e entrega de software."),
        ("Profissionais com base em desenvolvimento",
         "Desejam evoluir na carreira e ampliar repertório em práticas modernas de engenharia."),
    ]

    col_w = Inches(3.9)
    row_h = Inches(2.0)
    positions = [
        (Inches(0.3),  Inches(1.9)),
        (Inches(4.4),  Inches(1.9)),
        (Inches(8.5),  Inches(1.9)),
        (Inches(0.3),  Inches(4.1)),
        (Inches(4.4),  Inches(4.1)),
    ]
    for (left, top), (title, body) in zip(positions, cards):
        add_rect(slide, left, top, col_w, row_h, WHITE,
                 line_color=RGBColor(0xEC, 0xD6, 0xC2))
        add_textbox(slide, title, left + Inches(0.1), top + Inches(0.1),
                    col_w - Inches(0.2), Inches(0.4),
                    font_size=13, bold=True, color=DARK)
        add_textbox(slide, body, left + Inches(0.1), top + Inches(0.55),
                    col_w - Inches(0.2), Inches(1.3),
                    font_size=11, color=MUTED)


def slide_learning(prs):
    """Slide 3 – O que você vai aprender"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "O que você vai aprender",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=DARK)
    add_textbox(slide,
                "Durante o workshop, você terá contato prático com os pilares da "
                "automação e boas práticas de CI/CD.",
                Inches(0.5), Inches(1.1), Inches(12), Inches(0.6),
                font_size=14, color=MUTED)

    left_items = [
        "Conceitos fundamentais de DevOps e CI/CD",
        "Como configurar pipelines de automação com GitHub Actions",
        "Automação de testes e validação de código",
        "Dicas de carreira em DevOps e automação",
    ]
    right_items = [
        "Deploy automatizado com Azure DevOps",
        "Boas práticas para build e release automáticos",
        "Exemplos práticos e hands-on para aplicar imediatamente",
    ]

    col_w = Inches(5.8)
    col_h = Inches(4.8)
    for col, items in [(Inches(0.4), left_items), (Inches(6.8), right_items)]:
        add_rect(slide, col, Inches(1.9), col_w, col_h, WHITE,
                 line_color=RGBColor(0xEC, 0xD6, 0xC2))
        add_bullet_list(slide, items,
                        col + Inches(0.2), Inches(2.1),
                        col_w - Inches(0.4), col_h - Inches(0.4),
                        font_size=14, color=DARK,
                        bullet_color=SECONDARY)


def slide_instructor(prs):
    """Slide 4 – Com quem vou aprender"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "Com quem vou aprender",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=DARK)

    # Instructor card
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
             WHITE, line_color=RGBColor(0xEC, 0xD6, 0xC2))

    add_textbox(slide, "Marcel Levinspuhl",
                Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
                font_size=20, bold=True, color=SECONDARY)

    bio = (
        "Profissional de tecnologia com mais de 20 anos de experiência em "
        "desenvolvimento de software, engenharia e automação de ambientes. "
        "Atuou no mercado brasileiro em empresas de tecnologia e hoje trabalha "
        "no mercado irlandês com foco em práticas modernas de DevOps, automação e cloud.\n\n"
        "Com sólida experiência em Java e .NET, construiu uma carreira que une "
        "conhecimento técnico profundo, visão prática do ciclo completo de desenvolvimento "
        "de software e forte capacidade de transformar processos complexos em fluxos "
        "mais ágeis, organizados e eficientes.\n\n"
        "No workshop, ele compartilha essa experiência de forma aplicada, mostrando "
        "como ferramentas como GitHub e Azure DevOps podem ser usadas na prática para "
        "automatizar etapas essenciais do desenvolvimento."
    )
    add_textbox(slide, bio,
                Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.4),
                font_size=13, color=MUTED)


def slide_resources(prs):
    """Slide 5 – O que esperar do workshop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "O que esperar do workshop",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=DARK)

    resources = [
        ("Aula prática",    "Conteúdo direto ao ponto, com foco em execução."),
        ("Duração",         "1 hora de imersão orientada para aplicação real."),
        ("Formato",         "Evento on-line, ao vivo, via Teams."),
        ("Vagas",           "Limitadas para manter a dinâmica do encontro."),
        ("Data",            "Será definida em breve. Garanta sua vaga gratuita!"),
        ("Certificado",     "Emitido ao final do curso. Gratuito e válido internacionalmente."),
    ]

    card_w = Inches(3.9)
    card_h = Inches(1.8)
    positions = [
        (Inches(0.3),  Inches(1.2)),
        (Inches(4.7),  Inches(1.2)),
        (Inches(9.1),  Inches(1.2)),
        (Inches(0.3),  Inches(3.3)),
        (Inches(4.7),  Inches(3.3)),
        (Inches(9.1),  Inches(3.3)),
    ]
    for (left, top), (title, body) in zip(positions, resources):
        add_rect(slide, left, top, card_w, card_h, WHITE,
                 line_color=RGBColor(0xEC, 0xD6, 0xC2))
        add_textbox(slide, title, left + Inches(0.1), top + Inches(0.1),
                    card_w - Inches(0.2), Inches(0.4),
                    font_size=14, bold=True, color=DARK)
        add_textbox(slide, body, left + Inches(0.1), top + Inches(0.55),
                    card_w - Inches(0.2), Inches(1.1),
                    font_size=11, color=MUTED)


def slide_cta(prs):
    """Slide 6 – CTA final"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, TEAL_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "Garanta sua vaga no workshop",
                Inches(0.6), Inches(1.5), Inches(12), Inches(1.1),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Em apenas 1 hora, veja como estruturar um pipeline real de\n"
                "build, testes e deploy usando GitHub e Azure DevOps.",
                Inches(0.6), Inches(2.8), Inches(12), Inches(1.0),
                font_size=18, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)

    # CTA button
    add_rect(slide, Inches(4.0), Inches(4.2), Inches(5.3), Inches(0.7),
             ACCENT)
    add_textbox(slide, "Garantir minha vaga gratuitamente",
                Inches(4.05), Inches(4.25), Inches(5.2), Inches(0.6),
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "https://forms.gle/15zNEB56RkPeBTQY9",
                Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
                font_size=12, color=RGBColor(0xD5, 0xF0, 0xF5),
                align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def generate():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_audience(prs)
    slide_learning(prs)
    slide_instructor(prs)
    slide_resources(prs)
    slide_cta(prs)

    output = "workshop_github_azure_devops.pptx"
    prs.save(output)
    print(f"Presentation saved: {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    generate()
